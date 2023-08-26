from pptx import Presentation
from GPTExplainer import sendToChatGPT  # Import the sendToChatGPT function
import os


# The different layouts
TITLE = 0
TITLE_AND_CONTENT = 1
SECTION_HEADER = 2
TWO_CONTENT = 3
COMPARISON = 4
TITLE_ONLY = 5
BLANK = 6
CONTENT_WITH_CAPTION = 7
PICTURE_WITH_CAPTION = 8

'''
extract the text from the slide
'''



def extract_text(slide):
    text = ""
    for shape in slide.shapes:
        if shape != slide.shapes.title:
            text += get_shape_info(shape)
    return text


'''
get's the text from the specific shape
'''


def get_shape_info(shape):
    txt = ""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                txt += run.text
    return txt


'''
extract extra notes from the slide
'''


def extract_notes(slide):
    notes_txt = ""
    if slide.has_notes_slide:
        notes_slide = slide.notes_slide
        # Access the notes content
        notes_text = ""
        for shape in notes_slide.shapes:
            notes_text += get_shape_info(shape)
    return notes_txt



'''
gets extra information about the powerpoint from the useer
'''


def readExtraInfo(pPname):
    extra_info = {}
    extra_info['age'] = input("What age group is the PowerPoint designed for?")
    extra_info['lesson'] = input("What course is the PowerPoint for?")
    extra_info['subject'] = input("What is the subject of the lesson")
    extra_info['pPname'] = pPname
    extra_info['info'] = Add_Info_Slide()
    return extra_info



'''
returns a dictionary with extra information for certain slides
'''

def Add_Info_Slide():
    add_additional_info = ""
    while add_additional_info != "yes" and add_additional_info != "no":
        add_additional_info = input(
            "Do you want to add additional information for specific slides? (yes/no): ").lower().strip()
        if add_additional_info != "yes" and add_additional_info != "no":
            print("Invalid input. Please enter 'yes' or 'no'.")

    # Initialize a dictionary to store additional information for specific slides
    additional_info_dict = {}
    if add_additional_info == "yes":
        slide_numbers = input(
            "Enter the slide numbers for which you want to add additional information (separated by spaces): ")
        slide_numbers = slide_numbers.split()
        for slide_number in slide_numbers:
            slide_number = int(slide_number)
            additional_info = input(f"Please enter additional information for Slide {slide_number}: ")
            additional_info_dict[slide_number] = additional_info
    return additional_info_dict


'''
reads the powerpoint
'''
def read_powerpoint(pPname):
    # Load PowerPoint presentation
    #presentation = Presentation(pPname + '.pptx')
    presentation=Presentation(pPname)

    # Get the filename without the directory path
    filename = os.path.basename(pPname)

    # Split the filename based on underscores
    filename_parts = filename.split('_')

    # Extract the PowerPoint name
    pPname= " ".join(filename_parts[:-2])  # Join the parts except the last two

    # Create a list of dictionaries to store slide information
    slides = []

    extraInfo = readExtraInfo(pPname)
    # Initialize the explanation_data dictionary with the powerpoint_name field
    explanation_data = {
        'powerpoint_name': extraInfo['pPname']
    }

    print(presentation)

    # Extract slide information
    for slide in presentation.slides:
        slide_info = {}

        # Extract header
        slide_info['slide_header'] = slide.shapes.title.text if slide.shapes.title else ""

        # Extract text
        text=extract_text(slide)
        slide_info['slide_content'] = text
        if text=="":
            continue
        # extract notes
        slide_info['extra_notes'] = extract_notes(slide)


        additional_info_dict = extraInfo['info']
        # Add general information fields
        slide_info['age_group'] = extraInfo['age']
        slide_info['course'] = extraInfo['lesson']
        slide_info['subject_lesson'] = extraInfo['subject']
        slide_info['powerpoint_name']= extraInfo['pPname']

        # Add additional information if it exists for the current slide
        slide_number = presentation.slides.index(slide) + 1
        if slide_number in additional_info_dict:
            slide_info['additional_info'] = additional_info_dict[slide_number]
        else:
            slide_info['additional_info'] = ""

        # Append slide_info to the slides list
        slides.append(slide_info)

        # send the slide info to chatgpt
        explanation_data=sendToChatGPT(explanation_data,slide_info,slide_number)

    return explanation_data